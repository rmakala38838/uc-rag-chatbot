"""
Generate a professional 10-slide PowerPoint presentation for the UC RAG Chatbot project.
Slide 1: Title | Slides 2-9: Content | Slide 10: Thank You
Includes: architecture, tech stack, pipeline, retrieval, evaluation metrics, ethics/governance.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE

# UC Brand Colors
UC_NAVY = RGBColor(0x1B, 0x36, 0x5D)
UC_RED = RGBColor(0xC8, 0x10, 0x2E)
UC_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
UC_GRAY = RGBColor(0x33, 0x33, 0x33)
UC_DARK = RGBColor(0x00, 0x00, 0x00)
UC_LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
UC_LIGHT_BORDER = RGBColor(0xE2, 0xE5, 0xEA)
UC_GREEN = RGBColor(0x16, 0xA3, 0x4A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_slide_header(slide, title_text):
    """White header bar with red bottom border."""
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.95))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = UC_WHITE
    hdr.line.fill.background()
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.93), prs.slide_width, Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = UC_RED
    line.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.2), Inches(10), Inches(0.7))
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = UC_NAVY
    p.font.name = "Segoe UI"


def add_rounded_box(slide, left, top, width, height, fill_color, border_color=None, text="", font_size=11, font_color=UC_DARK, bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.alignment = align
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, color=UC_DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.alignment = align
    return tf


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

accent_top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
accent_top.fill.solid()
accent_top.fill.fore_color.rgb = UC_RED
accent_top.line.fill.background()

left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), prs.slide_height)
left_bar.fill.solid()
left_bar.fill.fore_color.rgb = UC_NAVY
left_bar.line.fill.background()

add_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
         "RAG-Powered AI Chatbot", font_size=42, color=UC_NAVY, bold=True)
add_text(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
         "University of the Cumberlands", font_size=28, color=UC_RED, bold=True)
add_text(slide, Inches(1.5), Inches(4.0), Inches(8), Inches(0.6),
         "Intelligent Q&A System  |  2,092 Web Pages  |  20,520 Knowledge Chunks", font_size=16, color=UC_GRAY)

bottom_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.0), Inches(3), Inches(0.04))
bottom_accent.fill.solid()
bottom_accent.fill.fore_color.rgb = UC_RED
bottom_accent.line.fill.background()

add_text(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.4),
         "Mounica Dayana  |  Eswari Ankitha Datla  |  Rajesh Makala  |  Sandeep Bagam", font_size=13, color=UC_DARK)
add_text(slide, Inches(1.5), Inches(5.6), Inches(8), Inches(0.4),
         "Ethics in Artificial Intelligence  |  University of the Cumberlands  |  July 2026", font_size=12, color=UC_GRAY)


# ============================================================
# SLIDE 2: Problem & Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "The Challenge & Our Solution")

add_rounded_box(slide, Inches(0.7), Inches(1.3), Inches(5.8), Inches(0.5),
                UC_NAVY, text="THE CHALLENGE", font_size=12, font_color=UC_WHITE, bold=True)
items = [
    "2,400+ pages of information across ucumberlands.edu",
    "Students struggle to find specific answers quickly",
    "Repetitive questions burden staff and admissions",
]
for i, item in enumerate(items):
    add_text(slide, Inches(1.0), Inches(2.0 + i * 0.5), Inches(5.5), Inches(0.45),
             f"•  {item}", font_size=13, color=UC_DARK)

add_rounded_box(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.5),
                UC_RED, text="OUR SOLUTION", font_size=12, font_color=UC_WHITE, bold=True)
solutions = [
    "AI chatbot with natural language understanding",
    "Instant answers grounded in university data",
    "Source citations with direct links to UC pages",
]
for i, item in enumerate(solutions):
    add_text(slide, Inches(7.3), Inches(2.0 + i * 0.5), Inches(5.5), Inches(0.45),
             f"•  {item}", font_size=13, color=UC_DARK)

# Metrics
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(3.9), Inches(12), Inches(0.02))
div.fill.solid()
div.fill.fore_color.rgb = UC_LIGHT_BORDER
div.line.fill.background()

metrics = [("2,092", "Pages Scraped"), ("20,520", "Vector Chunks"), ("768-dim", "Embeddings"), ("< 3 sec", "Response Time")]
for i, (val, label) in enumerate(metrics):
    left = Inches(1.0 + i * 3.1)
    add_text(slide, left, Inches(4.2), Inches(2.5), Inches(0.6),
             val, font_size=26, color=UC_RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left, Inches(4.75), Inches(2.5), Inches(0.4),
             label, font_size=11, color=UC_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3: System Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "System Architecture")

# INGESTION PIPELINE
add_text(slide, Inches(0.7), Inches(1.15), Inches(4), Inches(0.35),
         "DATA INGESTION PIPELINE", font_size=11, color=UC_RED, bold=True)

ing_y = Inches(1.55)
ing_h = Inches(1.1)
ing_w = Inches(2.2)

boxes_ing = [
    ("Web Scraper\n\n2,092 pages", Inches(0.7)),
    ("Markdown\nConverter\n\nJSON → .md", Inches(3.25)),
    ("Section-Aware\nChunking\n\n50–1500 chars", Inches(5.8)),
    ("Local\nEmbeddings\n\nBGE-base 768d", Inches(8.35)),
    ("ChromaDB\nVector Store\n\n20,520 chunks", Inches(10.9)),
]
for text, left in boxes_ing:
    add_rounded_box(slide, left, ing_y, ing_w, ing_h,
                    UC_LIGHT_BG, border_color=UC_NAVY, text=text, font_size=10, font_color=UC_NAVY)

arrow_y = ing_y + Inches(0.55)
for i in range(len(boxes_ing) - 1):
    cxn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
        boxes_ing[i][1] + ing_w + Inches(0.05), arrow_y,
        boxes_ing[i+1][1] - Inches(0.05), arrow_y)
    cxn.line.color.rgb = UC_NAVY
    cxn.line.width = Pt(1.5)

# QUERY PIPELINE
add_text(slide, Inches(0.7), Inches(3.1), Inches(4), Inches(0.35),
         "QUERY & RESPONSE PIPELINE", font_size=11, color=UC_RED, bold=True)

query_y = Inches(3.5)
boxes_query = [
    ("User\nQuestion\n\nChat Widget", Inches(0.7)),
    ("Query\nExpansion\n\nContext-aware", Inches(3.25)),
    ("Semantic\nSearch\n\nTop-15 → Top-10", Inches(5.8)),
    ("AWS Bedrock\nClaude Sonnet\n\nLLM Generation", Inches(8.35)),
    ("Formatted\nResponse\n\nLinks + Cite", Inches(10.9)),
]
for text, left in boxes_query:
    add_rounded_box(slide, left, query_y, ing_w, ing_h,
                    UC_LIGHT_BG, border_color=UC_RED, text=text, font_size=10, font_color=UC_DARK)

arrow_y2 = query_y + Inches(0.55)
for i in range(len(boxes_query) - 1):
    cxn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
        boxes_query[i][1] + ing_w + Inches(0.05), arrow_y2,
        boxes_query[i+1][1] - Inches(0.05), arrow_y2)
    cxn.line.color.rgb = UC_RED
    cxn.line.width = Pt(1.5)

# Vertical link
vert_x = Inches(6.9)
cxn_v = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, vert_x, ing_y + ing_h, vert_x, query_y)
cxn_v.line.color.rgb = UC_GRAY
cxn_v.line.width = Pt(1.5)
cxn_v.line.dash_style = 2

add_text(slide, Inches(0.7), Inches(5.1), Inches(12), Inches(0.4),
         "Python 3.11  •  FastAPI  •  sentence-transformers  •  ChromaDB (HNSW)  •  AWS Bedrock  •  HTML/CSS/JS",
         font_size=12, color=UC_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4: Technology Stack
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "Technology Stack")

col_data = [
    ("Backend", UC_NAVY, [
        "FastAPI + Uvicorn",
        "ChromaDB (cosine HNSW)",
        "AWS Bedrock Claude Sonnet",
        "sentence-transformers",
        "BAAI/bge-base-en-v1.5 (768d)",
    ]),
    ("Data Pipeline", UC_RED, [
        "Sitemap + link crawling",
        "JSON → Markdown converter",
        "Section-aware chunking",
        "Local GPU embeddings (MPS)",
        "4 min full ingestion",
    ]),
    ("Frontend", UC_NAVY, [
        "UC-branded responsive UI",
        "Floating chat widget",
        "Real-time typing indicators",
        "Clickable source links",
        "Context-aware follow-ups",
    ]),
]
for col_idx, (title, color, items) in enumerate(col_data):
    left = Inches(0.7 + col_idx * 4.2)
    add_rounded_box(slide, left, Inches(1.3), Inches(3.8), Inches(0.5),
                    color, text=title, font_size=13, font_color=UC_WHITE, bold=True)
    for i, item in enumerate(items):
        add_text(slide, left + Inches(0.2), Inches(2.0 + i * 0.45), Inches(3.6), Inches(0.42),
                 f"•  {item}", font_size=12, color=UC_DARK)

add_rounded_box(slide, Inches(0.7), Inches(4.6), Inches(12), Inches(0.55),
                UC_LIGHT_BG, border_color=UC_LIGHT_BORDER,
                text="Key Decision: Local embeddings (4 min) vs API-based (50 min) — 12x faster ingestion with same quality",
                font_size=12, font_color=UC_NAVY)


# ============================================================
# SLIDE 5: Data Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "Data Pipeline")

steps = [
    ("1", "SCRAPE", "Sitemap + link crawl → 2,092 structured JSON pages"),
    ("2", "CONVERT", "JSON → clean Markdown with YAML front matter"),
    ("3", "CHUNK", "Split by headings, sentence-boundary aware (50–1500 chars)"),
    ("4", "EMBED", "BAAI/bge-base-en-v1.5, 768 dims, Apple Silicon MPS GPU"),
    ("5", "STORE", "ChromaDB cosine HNSW index, persistent storage"),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.3 + i * 0.85)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.42), Inches(0.42))
    circle.fill.solid()
    circle.fill.fore_color.rgb = UC_RED
    circle.line.fill.background()
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(11)
    p.font.color.rgb = UC_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, Inches(1.5), y + Inches(0.02), Inches(1.5), Inches(0.35),
             title, font_size=12, color=UC_NAVY, bold=True)
    add_text(slide, Inches(3.0), y + Inches(0.02), Inches(5), Inches(0.4),
             desc, font_size=11, color=UC_DARK)

# Right side: key details
add_rounded_box(slide, Inches(8.5), Inches(1.3), Inches(4.3), Inches(0.45),
                UC_NAVY, text="CHUNKING STRATEGY", font_size=11, font_color=UC_WHITE, bold=True)
chunk_items = [
    "Split at markdown headings (##, ###)",
    "Context prefix: page title + section",
    "Sentence-boundary aware splits",
    "Min 50 / Max 1500 characters",
    "~10 chunks per page average",
]
for i, item in enumerate(chunk_items):
    add_text(slide, Inches(8.7), Inches(1.9 + i * 0.4), Inches(4.1), Inches(0.38),
             f"•  {item}", font_size=11, color=UC_DARK)


# ============================================================
# SLIDE 6: Retrieval & Grounding
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "Retrieval & Grounding Design")

# Left
add_rounded_box(slide, Inches(0.7), Inches(1.3), Inches(5.5), Inches(0.45),
                UC_NAVY, text="RETRIEVAL STRATEGY", font_size=11, font_color=UC_WHITE, bold=True)
ret_items = [
    "Query expansion resolves pronouns & follow-ups",
    "Same BGE model for indexing and querying",
    "Top-15 retrieved via cosine similarity",
    "Top-10 chunks passed as LLM context",
    "Relevance threshold: 0.3 minimum score",
    "Chat history (last 3 turns) maintained",
]
for i, item in enumerate(ret_items):
    add_text(slide, Inches(0.9), Inches(1.9 + i * 0.43), Inches(5.3), Inches(0.4),
             f"•  {item}", font_size=12, color=UC_DARK)

# Right
add_rounded_box(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(0.45),
                UC_RED, text="GROUNDING & REFUSAL", font_size=11, font_color=UC_WHITE, bold=True)
ground_items = [
    "Answers ONLY from retrieved university content",
    "No hallucinated or external information",
    "Every factual claim cited with source URL",
    "Low-confidence queries trigger refusal path",
    "Refusal template with contact info & links",
    "Out-of-domain questions detected and refused",
]
for i, item in enumerate(ground_items):
    add_text(slide, Inches(7.0), Inches(1.9 + i * 0.43), Inches(5.6), Inches(0.4),
             f"•  {item}", font_size=12, color=UC_DARK)

# Query expansion example
add_rounded_box(slide, Inches(0.7), Inches(4.7), Inches(12), Inches(0.6),
                UC_LIGHT_BG, border_color=UC_LIGHT_BORDER,
                text="Query Expansion Example:  \"tell me more\"  →  \"PhD in AI — tell me more\"  (prepends topic from chat history)",
                font_size=11, font_color=UC_NAVY)


# ============================================================
# SLIDE 7: UI/UX Design
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "UI/UX Design")

# Left
add_rounded_box(slide, Inches(0.7), Inches(1.3), Inches(5.5), Inches(0.45),
                UC_RED, text="BRAND IDENTITY", font_size=11, font_color=UC_WHITE, bold=True)
brand_items = [
    "UC official colors: Crimson #C8102E + Navy #1B365D",
    "White header with red accent border",
    "Red hero banner with white text",
    "Merriweather headings + Open Sans body",
    "Card-based topic browsing layout",
]
for i, item in enumerate(brand_items):
    add_text(slide, Inches(0.9), Inches(1.9 + i * 0.43), Inches(5.3), Inches(0.4),
             f"•  {item}", font_size=12, color=UC_DARK)

# Right
add_rounded_box(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(0.45),
                UC_NAVY, text="CHAT WIDGET", font_size=11, font_color=UC_WHITE, bold=True)
chat_items = [
    "Floating button (bottom-right, 68px)",
    "480px panel with navy header",
    "Typing indicator (3-dot animation)",
    "Navy user bubbles / white bot cards",
    "Clickable red links for citations",
]
for i, item in enumerate(chat_items):
    add_text(slide, Inches(7.0), Inches(1.9 + i * 0.43), Inches(5.6), Inches(0.4),
             f"•  {item}", font_size=12, color=UC_DARK)

# Bottom: Conversation rules
add_rounded_box(slide, Inches(0.7), Inches(4.4), Inches(12), Inches(0.45),
                UC_LIGHT_BG, border_color=UC_LIGHT_BORDER,
                text="Response Rules: Professional tone  •  No emojis  •  2-4 paragraphs  •  Clickable URLs  •  Never says 'based on context'",
                font_size=11, font_color=UC_NAVY)


# ============================================================
# SLIDE 8: Evaluation Metrics
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "Evaluation Results — 50 Golden Questions (Two Methods)")

# --- LEFT SIDE: Keyword-Based Evaluation ---
add_rounded_box(slide, Inches(0.5), Inches(1.15), Inches(6.0), Inches(0.4),
                UC_NAVY, text="METHOD 1: KEYWORD-BASED METRICS", font_size=10, font_color=UC_WHITE, bold=True)

kw_metrics = [
    ("89.8%", "P@5"),
    ("88.7%", "R@10"),
    ("73.6%", "Accuracy"),
    ("100%", "Faithful"),
    ("90%", "Refusal"),
]
for i, (val, label) in enumerate(kw_metrics):
    left = Inches(0.5 + i * 1.2)
    add_text(slide, left, Inches(1.65), Inches(1.15), Inches(0.45),
             val, font_size=18, color=UC_RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left, Inches(2.1), Inches(1.15), Inches(0.3),
             label, font_size=10, color=UC_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.5), Inches(2.5), Inches(6.0), Inches(0.25),
         "Automated keyword matching against expected terms per question",
         font_size=10, color=UC_GRAY)

# --- RIGHT SIDE: LLM-as-Judge Evaluation ---
add_rounded_box(slide, Inches(6.8), Inches(1.15), Inches(6.0), Inches(0.4),
                UC_RED, text="METHOD 2: LLM-AS-JUDGE (Claude Opus)", font_size=10, font_color=UC_WHITE, bold=True)

llm_metrics = [
    ("80.8%", "Overall"),
    ("87.2%", "Relevance"),
    ("82.8%", "Faithful"),
    ("76.4%", "Complete"),
    ("84.0%", "Citation"),
]
for i, (val, label) in enumerate(llm_metrics):
    left = Inches(6.8 + i * 1.2)
    add_text(slide, left, Inches(1.65), Inches(1.15), Inches(0.45),
             val, font_size=18, color=UC_RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left, Inches(2.1), Inches(1.15), Inches(0.3),
             label, font_size=10, color=UC_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(6.8), Inches(2.5), Inches(6.0), Inches(0.25),
         "Claude Opus evaluates each response vs ground-truth on 5 dimensions",
         font_size=10, color=UC_GRAY)

# --- BOTTOM: Per-Category (LLM Judge) ---
UC_BLACK = RGBColor(0x00, 0x00, 0x00)

div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.8), Inches(12.3), Inches(0.02))
div.fill.solid()
div.fill.fore_color.rgb = UC_LIGHT_BORDER
div.line.fill.background()

add_text(slide, Inches(0.5), Inches(3.0), Inches(6), Inches(0.3),
         "LLM-AS-JUDGE PER-CATEGORY (Claude Opus Evaluator)", font_size=11, color=UC_RED, bold=True)

cat_rows = [
    ("Category", "Correct", "Complete", "Faithful", "Relevant", "Citation"),
    ("Admissions (10)", "70%", "70%", "76%", "84%", "84%"),
    ("Academics (10)", "66%", "68%", "80%", "82%", "78%"),
    ("Tuition & Aid (8)", "85%", "88%", "88%", "98%", "95%"),
    ("Student Life (7)", "71%", "77%", "86%", "91%", "89%"),
    ("General (10)", "76%", "78%", "86%", "90%", "78%"),
    ("Out-of-Domain (5)", "76%", "84%", "84%", "76%", "84%"),
]
for i, row in enumerate(cat_rows):
    y = Inches(3.4 + i * 0.33)
    is_hdr = (i == 0)
    for j, cell in enumerate(row):
        left = Inches(0.5 + j * 1.3)
        add_text(slide, left, y, Inches(1.25), Inches(0.3),
                 cell, font_size=11, color=UC_NAVY if is_hdr else UC_BLACK, bold=is_hdr)

# --- BOTTOM RIGHT: Performance & Method ---
add_text(slide, Inches(8.5), Inches(3.0), Inches(4.3), Inches(0.3),
         "PERFORMANCE & METHOD", font_size=11, color=UC_RED, bold=True)
perf = [
    "Source retrieval rate: 68.9%",
    "Avg similarity: 0.7109",
    "Avg retrieval: 0.079 sec",
    "Avg generation: 3.7 sec",
    "50/50 tests, 0 errors",
    "Ground-truth from UC docs",
    "Judge: Claude Opus via Bedrock",
]
for i, item in enumerate(perf):
    add_text(slide, Inches(8.5), Inches(3.4 + i * 0.33), Inches(4.3), Inches(0.3),
             f"•  {item}", font_size=11, color=UC_BLACK)


# ============================================================
# SLIDE 9: Ethics, Safety & Governance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "Ethics, Safety & Governance")

# Left column: Crawl Etiquette
add_rounded_box(slide, Inches(0.7), Inches(1.2), Inches(3.8), Inches(0.42),
                UC_NAVY, text="CRAWL ETIQUETTE", font_size=10, font_color=UC_WHITE, bold=True)
crawl_items = [
    "100ms delay between requests",
    "Respects robots.txt directives",
    "Only public pages (no login-gated)",
    "Identified bot User-Agent header",
    "Data for academic use only",
]
for i, item in enumerate(crawl_items):
    add_text(slide, Inches(0.9), Inches(1.75 + i * 0.35), Inches(3.6), Inches(0.33),
             f"•  {item}", font_size=10, color=UC_DARK)

# Middle: PII/PHI
add_rounded_box(slide, Inches(4.8), Inches(1.2), Inches(3.8), Inches(0.42),
                UC_RED, text="PII / PHI EXCLUSION", font_size=10, font_color=UC_WHITE, bold=True)
pii_items = [
    "No student personal data collected",
    "Faculty info is publicly published",
    "No login-protected content",
    "User queries not stored or logged",
    "No health information (PHI) accessed",
]
for i, item in enumerate(pii_items):
    add_text(slide, Inches(5.0), Inches(1.75 + i * 0.35), Inches(3.6), Inches(0.33),
             f"•  {item}", font_size=10, color=UC_DARK)

# Right: Safety
add_rounded_box(slide, Inches(8.9), Inches(1.2), Inches(3.8), Inches(0.42),
                UC_NAVY, text="AI SAFETY GUARDRAILS", font_size=10, font_color=UC_WHITE, bold=True)
safety_items = [
    "Only answers from retrieved context",
    "Relevance threshold for refusal",
    "No hallucinated information",
    "Source URLs cited in every answer",
    "100% faithfulness in evaluation",
]
for i, item in enumerate(safety_items):
    add_text(slide, Inches(9.1), Inches(1.75 + i * 0.35), Inches(3.6), Inches(0.33),
             f"•  {item}", font_size=10, color=UC_DARK)

# Bottom: Risk table
add_text(slide, Inches(0.7), Inches(3.7), Inches(5), Inches(0.3),
         "ETHICAL RISKS & MITIGATIONS", font_size=10, color=UC_RED, bold=True)

risks = [
    ("Risk", "Mitigation Strategy"),
    ("Outdated information", "Timestamps; disclaimers; scheduled re-scraping pipeline"),
    ("Incorrect academic advice", "Refusal on low confidence; links to official pages; no guarantees"),
    ("Over-reliance on AI", "Footer disclaimer: 'For official inquiries, contact UC directly'"),
    ("Bias in embeddings", "Open-source BGE model; tested across 7 categories for fairness"),
    ("Privacy of user queries", "No query logging; session-only history; no PII stored"),
]
for i, (risk, mitigation) in enumerate(risks):
    y = Inches(4.0 + i * 0.3)
    is_hdr = (i == 0)
    add_text(slide, Inches(0.7), y, Inches(3.2), Inches(0.28),
             risk, font_size=9, color=UC_NAVY if is_hdr else UC_DARK, bold=is_hdr)
    add_text(slide, Inches(4.0), y, Inches(9), Inches(0.28),
             mitigation, font_size=9, color=UC_NAVY if is_hdr else UC_DARK, bold=is_hdr)


# ============================================================
# SLIDE 10: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

accent_top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
accent_top.fill.solid()
accent_top.fill.fore_color.rgb = UC_RED
accent_top.line.fill.background()

left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), prs.slide_height)
left_bar.fill.solid()
left_bar.fill.fore_color.rgb = UC_NAVY
left_bar.line.fill.background()

add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1),
         "Thank You", font_size=44, color=UC_NAVY, bold=True)

add_text(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.6),
         "Questions & Live Demo", font_size=24, color=UC_RED, bold=True)

div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.2), Inches(4), Inches(0.03))
div.fill.solid()
div.fill.fore_color.rgb = UC_RED
div.line.fill.background()

add_text(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.5),
         "Demo: localhost:8080", font_size=14, color=UC_GRAY)
add_text(slide, Inches(1.5), Inches(5.1), Inches(10), Inches(0.5),
         "GitHub: github.com/rmakala38838/uc-rag-chatbot", font_size=14, color=UC_GRAY)
add_text(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.4),
         "Mounica Dayana  |  Eswari Ankitha Datla  |  Rajesh Makala  |  Sandeep Bagam", font_size=13, color=UC_DARK)
add_text(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.4),
         "Ethics in Artificial Intelligence  |  University of the Cumberlands  |  July 2026", font_size=12, color=UC_GRAY)


# Save
output_path = "UC_RAG_Chatbot_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Slides: {len(prs.slides)}")